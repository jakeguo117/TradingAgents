"""PowerPoint presentation generator — Sun Life brand style."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .gemini_client import generate_structured
from .prompts import PPTX_PROMPT
from .schemas import SlideContent

# Sun Life brand colors
TEAL_DARK = RGBColor(0x00, 0x39, 0x46)    # #003946
GOLD = RGBColor(0xFE, 0xBE, 0x10)         # #FEBE10
BODY_TEXT = RGBColor(0x22, 0x22, 0x22)     # #222222
GRAY_TEXT = RGBColor(0x5B, 0x5E, 0x5E)     # #5B5E5E
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY_BG = RGBColor(0xF5, 0xF5, 0xF5)


def _apply_slide_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def generate_pptx(
    report_text: str, output_dir: Path, model: str = "gemini-2.5-flash"
) -> Path:
    """Generate a Sun Life-styled PowerPoint deck from a trading analysis report.

    Args:
        report_text: The complete trading analysis report markdown.
        output_dir: Directory to save the presentation.
        model: Gemini model to use.

    Returns:
        Path to the generated .pptx file.
    """
    slides_data: list[SlideContent] = generate_structured(
        report_text=report_text,
        prompt=PPTX_PROMPT,
        response_schema=list[SlideContent],
        model=model,
    )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides_data):
        if i == 0:
            # Title slide — dark teal background
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            _apply_slide_background(slide, TEAL_DARK)

            # Title text box
            from pptx.util import Emu
            title_box = slide.shapes.add_textbox(
                Inches(1), Inches(2), Inches(11), Inches(2)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data["title"]
            p.font.size = Pt(40)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.LEFT

            # Subtitle / bullets
            for bullet in slide_data["bullets"]:
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(20)
                p.font.color.rgb = GOLD
                p.alignment = PP_ALIGN.LEFT

            # Gold accent bar
            bar = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(1), Inches(1.7), Inches(2), Inches(0.06)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = GOLD
            bar.line.fill.background()

        else:
            # Content slides — white background
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

            # Teal header bar
            header_bar = slide.shapes.add_shape(
                1, Inches(0), Inches(0), prs.slide_width, Inches(1.2)
            )
            header_bar.fill.solid()
            header_bar.fill.fore_color.rgb = TEAL_DARK
            header_bar.line.fill.background()

            # Title on header
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.2), Inches(11), Inches(0.8)
            )
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data["title"]
            p.font.size = Pt(28)
            p.font.color.rgb = WHITE
            p.font.bold = True

            # Gold accent line under header
            accent = slide.shapes.add_shape(
                1, Inches(0.8), Inches(1.2), Inches(2), Inches(0.05)
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = GOLD
            accent.line.fill.background()

            # Bullet content
            content_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(1.6), Inches(11), Inches(5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True

            for j, bullet in enumerate(slide_data["bullets"]):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"  {bullet}"
                p.font.size = Pt(18)
                p.font.color.rgb = BODY_TEXT
                p.space_after = Pt(12)

                # Gold bullet marker prefix
                if p.runs:
                    p.runs[0].font.color.rgb = BODY_TEXT

        # Speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = slide_data.get("speaker_notes", "")

    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / "trading_report.pptx"
    prs.save(str(output_path))

    return output_path
